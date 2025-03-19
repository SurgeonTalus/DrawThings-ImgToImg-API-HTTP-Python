import os
import json
import base64
import requests
from pptx import Presentation

# Set paths
DOWNLOADS_FOLDER = os.path.expanduser("~/Downloads")
PPTX_FILE = os.path.join(DOWNLOADS_FOLDER, "powerpoint.pptx")

# API Endpoints
LM_STUDIO_URL = "http://localhost:1234/v1/chat/completions"
DRAW_THINGS_URL = "http://127.0.0.1:7860/sdapi/v1/img2img"

# Image generation settings
STEPS = 4
REMOVE_BACKGROUND = False  # Change to True if using background removal

# Function to extract text from slides
def extract_text_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    slides_text = []
    
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text.strip())
        slides_text.append("\n".join(slide_text))
    
    return slides_text

# Function to get an image description from LM Studio
def description_prompt_text(prompt_text):
    request_data = {
        "messages": [
            {"role": "system", "content": "Do not chat. Answer precisely. Using the context of the PowerPoint slide, come up with a description of a fitting image."},
            {"role": "user", "content": prompt_text}
        ],
        "model": "gemma-3-12b-it",
        "temperature": 0.1,
        "max_tokens": 500,
        "stream": False,
        "reset_history": True,
    }

    response = requests.post(LM_STUDIO_URL, json=request_data)
    
    if response.status_code == 200:
        response_json = response.json()
        return response_json.get("choices", [{}])[0].get("message", {}).get("content", "No description generated.")
    else:
        print(f"Error generating description. Status code: {response.status_code}")
        return None

# Function to generate an image using Draw Things API (img2img support)
def generate_image(prompt, base64_image=None):
    print(f"Generating image with {STEPS} steps...")

    params = {
        "prompt": prompt,
        "negative_prompt": "(bokeh, worst quality, low quality, normal quality, (variations):1.4), blur:1.5",
        "seed": 4068245935,
        "steps": STEPS,
        "guidance_scale": 10,
        "batch_count": 1
    }

    # If base64_image is provided, use img2img mode
    if base64_image:
        params["init_images"] = [base64_image]
        params["denoising_strength"] = 0.75  # Adjust for img2img effect strength

    headers = {"Content-Type": "application/json"}
    response = requests.post(DRAW_THINGS_URL, json=params, headers=headers)

    if response.status_code == 200:
        data = response.json()
        images = data.get("images", [])
        if images:
            temp_image_path = os.path.join("/tmp", "generated_image.png")
            with open(temp_image_path, "wb") as img_file:
                img_file.write(base64.b64decode(images[0]))
            return temp_image_path

    print(f"Error generating image: {response.status_code}, {response.text}")
    return None

# Function to read an image and convert to base64
def encode_image_to_base64(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")

# Main script execution
if __name__ == "__main__":
    if not os.path.exists(PPTX_FILE):
        print(f"Error: PowerPoint file not found at {PPTX_FILE}")
        exit(1)

    slides = extract_text_from_pptx(PPTX_FILE)

    for i, slide_text in enumerate(slides, start=1):
        print(f"\nProcessing Slide {i}:")
        print(slide_text)

        # Get description from LM Studio
        description = description_prompt_text(slide_text)
        if not description:
            print(f"Skipping slide {i} due to missing description.")
            continue

        print(f"Generated Description: {description}")

        # (Optional) Load base64 image if needed
        base64_image = None
        image_path = os.path.join(DOWNLOADS_FOLDER, f"slide_{i}_image.png")
        if os.path.exists(image_path):  # Use existing image if available
            base64_image = encode_image_to_base64(image_path)

        # Generate image using Draw Things API
        generated_image_path = generate_image(description, base64_image=base64_image)
        if generated_image_path:
            print(f"Image for Slide {i} saved at: {generated_image_path}")
        else:
            print(f"Failed to generate image for Slide {i}.")